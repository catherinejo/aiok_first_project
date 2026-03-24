"""File parsing utilities for PPT, PDF, and Word documents."""

from __future__ import annotations

import io
from pathlib import Path

import pdfplumber
from docx import Document
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt"}


def parse_pdf(file_content: bytes) -> str:
    """Extract text from PDF file."""
    text_parts = []
    with pdfplumber.open(io.BytesIO(file_content)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"[Page {i}]\n{page_text}")
    return "\n\n".join(text_parts)


def parse_docx(file_content: bytes) -> str:
    """Extract text from Word document."""
    doc = Document(io.BytesIO(file_content))
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text_parts.append(row_text)
    return "\n".join(text_parts)


def parse_pptx(file_content: bytes) -> str:
    """Extract text from PowerPoint presentation."""
    prs = Presentation(io.BytesIO(file_content))
    text_parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        if slide_texts:
            text_parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(text_parts)


def parse_file(filename: str, content: bytes) -> str:
    """Parse file based on extension and return extracted text."""
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return parse_pdf(content)
    elif ext in (".docx", ".doc"):
        return parse_docx(content)
    elif ext in (".pptx", ".ppt"):
        return parse_pptx(content)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def is_supported_file(filename: str) -> bool:
    """Check if file type is supported."""
    ext = Path(filename).suffix.lower()
    return ext in SUPPORTED_EXTENSIONS
